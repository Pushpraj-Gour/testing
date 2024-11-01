from pptx import Presentation
from pptx.util import Pt

# Function to add content to a slide with dynamic font size based on the level of nesting
def add_content_to_slide(slide, content, level=0):
    # Define a placeholder for the content
    content_placeholder = slide.shapes.placeholders[1]  # Assuming layout with Title and Content
    text_frame = content_placeholder.text_frame

    # Font size reduction logic (starting at 24 points and decreasing by 2 for each level)
    base_font_size = 24
    font_size = base_font_size - (level * 2)
    if font_size < 10:
        font_size = 10  # Set a minimum font size to maintain readability

    # If content is a list, handle each item
    if isinstance(content, list):
        for item in content:
            # Recursively add list items as content
            add_content_to_slide(slide, item, level + 1)
    else:
        # Add text to the slide with dynamic font size based on level
        p = text_frame.add_paragraph()
        p.text = str(content)  # Convert content to string if it's not already a string
        p.font.size = Pt(font_size)

# Recursive function to handle nested dictionaries and add content to slides
def handle_nested_dict(prs, title, nested_dict, level=0):
    slide_layout = prs.slide_layouts[1]  # 1 = Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the slide title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    # Iterate through the nested dictionary and add content
    for key, value in nested_dict.items():
        if isinstance(value, dict):
            # Recursive call for nested dictionaries
            add_content_to_slide(slide, key, level)  # Add key as content
            handle_nested_dict(prs, key, value, level + 1)  # Recursively handle nested dicts
        else:
            # Add key and value for non-nested content
            add_content_to_slide(slide, f"{key}: {value}", level)

# Create a PowerPoint presentation object
prs = Presentation()

# Sample nested data structure
data = {
    "Title Slide": {
        "Introduction": "This is the introduction to the presentation."
    },
    "Slide 1": {
        "Section 1": {
            "Subsection 1": "Details about subsection 1",
            "Subsection 2": "Details about subsection 2"
        },
        "Section 2": "Details about section 2"
    },
    "Slide 2": [
        {"List Item 1": "Details for list item 1"},
        {"List Item 2": "Details for list item 2"}
    ],
    "Final Slide": "This is the conclusion of the presentation."
}

# Iterate through the data and generate slides
for title, content in data.items():
    if isinstance(content, dict):
        handle_nested_dict(prs, title, content)
    else:
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        add_content_to_slide(slide, content)


# Save the PowerPoint presentation
prs.save('output.pptx')

print("Done")
