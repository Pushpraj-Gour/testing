""" Ref for slide types: 
0 -> title and subtitle 
1 -> title and content 
2 -> section header 
3 -> two content 
4 -> Comparison 
5 -> Title only 
6 -> Blank 
7 -> Content with caption 
8 -> Pic with caption 
"""


from pptx import Presentation

presentation = Presentation()

# Add a title slide
slide_1_layout = presentation.slide_layouts[0]  # Title Slide layout
slide_1 = presentation.slides.add_slide(slide_1_layout)

# Add title and subtitle to the title slide
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Welcome to Python-PPTX"
subtitle.text = "Creating PowerPoint Presentations with Python"

# Add a slide with a title and content
slide_2_layout = presentation.slide_layouts[1]  # Title and Content layout
slide_2 = presentation.slides.add_slide(slide_2_layout)

# Add title and content
slide_2_title = slide_2.shapes.title
slide_2_title.text = "Slide with Bulleted List"

# Add bullet points
content = slide_2.shapes.placeholders[1]
text_frame = content.text_frame # type: ignore

# First bullet point
text_frame.text = "This is the first bullet point"

# Additional bullet points
p = text_frame.add_paragraph()
p.text = "This is the second bullet point"

p = text_frame.add_paragraph()
p.text = "This is the third bullet point"

# Save the presentation
presentation.save("sample_presentation.pptx")

print("Presentation created successfully!")
