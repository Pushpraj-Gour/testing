from pptx import Presentation
from pptx.util import Pt

def create_ppt_from_data(data, pptx_file_name):

    presentation = Presentation()
    
    def add_slide(pres, title, content_text=None,level=0):

        slide_layout = pres.slide_layouts[1]  # Title and Content layout
        slide = pres.slides.add_slide(slide_layout)
        
        # Add title
        slide_title = slide.shapes.title
        slide_title.text = title
        

            # Create a single paragraph for the content

        base_font_size = 24
        font_size = base_font_size - (level * 2)
        if font_size < 10:
            font_size = 10 
            
    
            # paragraph = text_frame.add_paragraph() 
            # paragraph.text = str(content)  # Convert content to string if it's not already a string
            # paragraph.font.size = Pt(font_size)

            # paragraph  = text_frame.add_paragraph()
            # add_content_to_slide(paragraph, content)
        if content_text:
            content = slide.shapes.placeholders[1]
            text_frame = content.text_frame
            text_frame.text = content_text

        add_content_to_slide(text_frame, content,level+1)

    def add_content_to_slide(text_frame, content,level):

        if isinstance(content, str):
        
            text_frame.text = content + "\n" 
        elif isinstance(content, dict):
    
            for key, value in content.items():
           
                text_frame.text  key + "\n" 
                
                # Recursively handle nested content
                add_content_to_slide(text_frame, value,level+1)
        elif isinstance(content, list):
           
            for item in content:
                if isinstance(item, dict):
               
                    for key, value in item.items():
                        text_frame.text = key + "\n" 
                        add_content_to_slide(text_frame, value,level+1)
                else:
            
                    text_frame.text = str(item) + "\n" 

    for section_title, section_content in data.items():
    
        add_slide(presentation, section_title, section_content)


    presentation.save(pptx_file_name)
    print(f"Presentation '{pptx_file_name}' created successfully!")

data = {
    "Title Slide": {
        "Subtitle": "This is an example of a generic data-driven PPT",
    },
    "Section 1": {
        "Introduction": "This is an introductory section",
        "Details": {
            "Point 1": "Description of point 1",
            "Point 2": "Description of point 2",
            "Nested List": ["Item A", "Item B", "Item C"]
        }
    },
    "Section 2": {
        "Summary": "This section contains a summary.",
        "Points": [
            {"Subsection 1": "Subsection 1 details"},
            {"Subsection 2": "Subsection 2 details"}
        ]
    },
    "Conclusion": "This is the conclusion of the presentation."
}

create_ppt_from_data(data, "output.pptx")
