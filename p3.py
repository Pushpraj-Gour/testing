from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(presentation, title_text, content_text=None):

    slide_layout = presentation.slide_layouts[1]  # Title and Content layout
    slide = presentation.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text

    if content_text:
        content = slide.shapes.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = content_text

def process_dict(presentation, data, level=0):

    for key, value in data.items():
        if isinstance(value, dict):
            add_slide(presentation, f"{' ' * level * 2}{key}")
            process_dict(presentation, value, level + 1)
        elif isinstance(value, list):
            add_slide(presentation, f"{' ' * level * 2}{key}", ", ".join(map(str, value)))
        else:
            add_slide(presentation, f"{' ' * level * 2}{key}", str(value))

def process_list(presentation, data, level=0):

    for idx, item in enumerate(data):
        if isinstance(item, dict):

            add_slide(presentation, f"Item {idx + 1}")
            process_dict(presentation, item, level + 1) 

        elif isinstance(item, list):

            add_slide(presentation, f"List {idx + 1}", ", ".join(map(str, item)))
        else:

            add_slide(presentation, f"Item {idx + 1}", str(item))

def create_pptx_from_data(data, output_file="output_presentation.pptx"):

    presentation = Presentation()

    if isinstance(data, dict):
        process_dict(presentation, data)

    elif isinstance(data, list):
        process_list(presentation, data)

    else:
        # If the data is not a list or dictionary, just add it as a slide
        add_slide(presentation, "Data", str(data))

    # Save the presentation to the specified file
    presentation.save(output_file)
    print(f"Presentation saved as {output_file}")

# Example usage:
data = {
    "Title Slide": {
        "Subtitle": "This is an example of a generic data-driven PPT",
    },
    "Section 1": {
        "Introduction": "This is an introductory section",
        "Details": {
            "Point 1": "Description of point 1",
            "Point 2": "Description of point 2",
            "Nested List": ["Item A", "Item B", "Item C"]
        }
    },
    "Section 2": {
        "Summary": "This section contains a summary.",
        "Points": [
            {"Subsection 1": "Subsection 1 details"},
            {"Subsection 2": "Subsection 2 details"}
        ]
    },
    "Conclusion": "This is the conclusion of the presentation."
}

for page in data.items():
    


# Generate the presentation from the example data
create_pptx_from_data(data, "example_presentation.pptx")
