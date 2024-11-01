from pptx import Presentation
from pptx.util import Inches

data = {
    "Title Slide": {
        "Introduction": "This is the introduction to the presentation."
    },
    "Slide 1": {
        "Section 1": {
            "Subsection 1": "Details about subsection 1",
            "Subsection 2": "Details about subsection 2"
        },
        "Section 2": "Details about section 2"
    },
    "Slide 2": [
        {"List Item 1": "Details for list item 1"},
        {"List Item 2": "Details for list item 2"}
    ],
    "Final Slide": "This is the conclusion of the presentation."
}

def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  
    slide = prs.slides.add_slide(slide_layout)
    
    # Set the title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    # Set the content
    content_placeholder = slide.shapes.placeholders[1]
    if isinstance(content, dict):
        # Handle nested dictionaries
        add_nested_content(content_placeholder, content)
    elif isinstance(content, list):
        # Handle list of dictionaries or simple lists
        for item in content:
            if isinstance(item, dict):
                for key, value in item.items():
                    content_placeholder.text += f"{key}: {value}\n"
            else:
                content_placeholder.text += f"- {item}\n"
    else:
        # Handle simple text content
        content_placeholder.text = str(content)

# Recursive function to handle nested dictionaries
def add_nested_content(placeholder, nested_dict, level=0):
    indent = " " * (level * 2)  # Indentation for nested items
    for key, value in nested_dict.items():
        if isinstance(value, dict):
            placeholder.text += f"{indent}{key}:\n"
            add_nested_content(placeholder, value, level + 1)
        else:
            placeholder.text += f"{indent}{key}: {value}\n"

# Create a PowerPoint presentation
prs = Presentation()

# Iterate through the data structure and add slides
for title, content in data.items():
    add_slide(prs, title, content)

# Save the presentation
prs.save('nested_structure_presentation.pptx')

print("Presentation created successfully!")
