from pptx import Presentation
from pptx.util import Pt

def create_ppt_from_data(data, pptx_file_name):

    presentation = Presentation()
    
    def add_slide(pres, title, content=None):

        slide_layout = pres.slide_layouts[1]  # Title and Content layout
        slide = pres.slides.add_slide(slide_layout)
        
        slide_title = slide.shapes.title
        slide_title.text = title
        
        if content:
            content_placeholder = slide.shapes.placeholders[1]  # Placeholder for content
            text_frame = content_placeholder.text_frame
            paragraph  = text_frame.add_paragraph()
            add_content_to_slide(paragraph, content)

    def add_content_to_slide(paragraph, content):

        if isinstance(content, str):
            paragraph.text = content
        elif isinstance(content, dict):

            for key, value in content.items():  
                paragraph.text = key
                
                run = paragraph.runs[0]
                run.font.size = Pt(16)  # Set font size for main points
                
                # Recursively handle nested content
                add_content_to_slide(paragraph, value)

        elif isinstance(content, list):
        
            for item in content:
                if isinstance(item, dict):
                    # If the item in the list is a dictionary, recursively add its content
                    for key, value in item.items():

                        paragraph.text = key
                        run = paragraph.runs[0]
                        run.font.size = Pt(14)  # Set font size for list items
                        add_content_to_slide(paragraph, value)
                else:

                    paragraph.text = str(item)
                    run = paragraph.runs[0]
                    run.font.size = Pt(14)  # Set font size for list items

    for section_title, section_content in data.items():
        add_slide(presentation, section_title, section_content)

    # Save the presentation
    presentation.save(pptx_file_name)
    print(f"Presentation '{pptx_file_name}' created successfully!")

# Example data structure
data = {
    "Title Slide": {
        "Subtitle": "This is an testing ppt",
    },
    "Section 1": {
        "Introduction": "This is an introductory section",
        "Details": {
            "Point 1": "Description of point 1",
            "Point 2": "Description of point 2",
            "Nested List": ["Item A", "Item B", "Item C"]
        }
    },
    "Section 2": {
        "Summary": "This section contains a summary.",
        "Points": [
            {"Subsection 1": "Subsection 1 details"},
            {"Subsection 2": "Subsection 2 details"}
        ]
    },
    "Conclusion": "This is the conclusion of the presentation."
}

# Call the function with the data
create_ppt_from_data(data, "output.pptx")
